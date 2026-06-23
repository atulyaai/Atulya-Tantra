"""Document generation with optional rich renderers and safe fallbacks."""
from __future__ import annotations

import csv
import html
import re
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class DocumentContent:
    title: str
    summary: str = ""
    sections: list[dict[str, Any]] = field(default_factory=list)
    rows: list[dict[str, Any]] = field(default_factory=list)


@dataclass
class DocumentResult:
    path: str
    format: str
    fallback: bool = False
    metadata: dict[str, Any] = field(default_factory=dict)


def _slug(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9]+", "-", value.lower()).strip("-")[:80] or "document"


def content_from_prompt(prompt: str) -> DocumentContent:
    lines = [line.strip(" -\t") for line in prompt.splitlines() if line.strip()]
    title = (lines[0] if lines else "Generated Document")[:90]
    body = " ".join(lines[1:] or [prompt]).strip()
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", body) if s.strip()]
    sentences = sentences or [prompt.strip()]
    return DocumentContent(
        title=title,
        summary=sentences[0],
        sections=[
            {"heading": "Overview", "body": sentences[0]},
            {"heading": "Key Points", "bullets": sentences[1:6] or [sentences[0]]},
            {"heading": "Next Steps", "bullets": ["Review the draft", "Add source data", "Finalize output"]},
        ],
    )


class DocumentEngine:
    """Generate PDF, DOCX, PPTX, XLSX, HTML, CSV, or Markdown."""

    def __init__(self, output_dir: str | Path = "assets/creations/documents"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def create_from_prompt(self, prompt: str, format: str = "markdown") -> DocumentResult:
        content = content_from_prompt(prompt)
        fmt = format.lower().lstrip(".")
        if fmt == "pdf":
            return self.create_pdf(content)
        if fmt == "docx":
            return self.create_docx(content)
        if fmt == "pptx":
            return self.create_pptx(content)
        if fmt == "xlsx":
            return self.create_xlsx(self._rows(content))
        if fmt == "html":
            return self.create_html(content)
        if fmt == "csv":
            return self.create_csv(self._rows(content))
        return self.create_markdown(content)

    def create_markdown(self, content: DocumentContent) -> DocumentResult:
        path = self._path(content.title, "md")
        chunks = [f"# {content.title}", "", content.summary]
        for section in content.sections:
            chunks.extend(["", f"## {section.get('heading', 'Section')}"])
            if section.get("body"):
                chunks.append(str(section["body"]))
            chunks.extend(f"- {item}" for item in section.get("bullets", []))
        path.write_text("\n".join(chunks).strip() + "\n", encoding="utf-8")
        return DocumentResult(str(path), "markdown")

    def create_html(self, content: DocumentContent) -> DocumentResult:
        path = self._path(content.title, "html")
        parts = ["<!doctype html><meta charset='utf-8'>", f"<h1>{html.escape(content.title)}</h1>"]
        parts.append(f"<p>{html.escape(content.summary)}</p>")
        for section in content.sections:
            parts.append(f"<h2>{html.escape(str(section.get('heading', 'Section')))}</h2>")
            if section.get("body"):
                parts.append(f"<p>{html.escape(str(section['body']))}</p>")
            if section.get("bullets"):
                parts.append("<ul>" + "".join(f"<li>{html.escape(str(x))}</li>" for x in section["bullets"]) + "</ul>")
        path.write_text("\n".join(parts), encoding="utf-8")
        return DocumentResult(str(path), "html")

    def create_pdf(self, content: DocumentContent) -> DocumentResult:
        try:
            from fpdf import FPDF
            path = self._path(content.title, "pdf")
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 18)
            pdf.multi_cell(0, 10, content.title)
            pdf.set_font("Helvetica", size=11)
            pdf.multi_cell(0, 8, content.summary)
            for section in content.sections:
                pdf.set_font("Helvetica", "B", 13)
                pdf.multi_cell(0, 8, str(section.get("heading", "Section")))
                pdf.set_font("Helvetica", size=11)
                pdf.multi_cell(0, 7, str(section.get("body", "")))
                for item in section.get("bullets", []):
                    pdf.multi_cell(0, 7, f"- {item}")
            pdf.output(str(path))
            return DocumentResult(str(path), "pdf")
        except Exception as exc:
            return self._fallback(content, "pdf", exc)

    def create_docx(self, content: DocumentContent) -> DocumentResult:
        try:
            from docx import Document
            path = self._path(content.title, "docx")
            doc = Document()
            doc.add_heading(content.title, 0)
            doc.add_paragraph(content.summary)
            for section in content.sections:
                doc.add_heading(str(section.get("heading", "Section")), level=1)
                if section.get("body"):
                    doc.add_paragraph(str(section["body"]))
                for item in section.get("bullets", []):
                    doc.add_paragraph(str(item), style="List Bullet")
            doc.save(path)
            return DocumentResult(str(path), "docx")
        except Exception as exc:
            return self._fallback(content, "docx", exc)

    def create_pptx(self, content: DocumentContent) -> DocumentResult:
        try:
            from pptx import Presentation
            path = self._path(content.title, "pptx")
            deck = Presentation()
            for title, bullets in [(content.title, [content.summary])] + [
                (str(s.get("heading", "Section")), list(s.get("bullets", [])) or [str(s.get("body", ""))])
                for s in content.sections
            ]:
                slide = deck.slides.add_slide(deck.slide_layouts[1])
                slide.shapes.title.text = title
                slide.placeholders[1].text = "\n".join(str(x) for x in bullets if x)
            deck.save(path)
            return DocumentResult(str(path), "pptx")
        except Exception as exc:
            return self._fallback(content, "pptx", exc)

    def create_xlsx(self, rows: list[dict[str, Any]]) -> DocumentResult:
        try:
            from openpyxl import Workbook
            path = self._path("spreadsheet", "xlsx")
            workbook = Workbook()
            sheet = workbook.active
            headers = list(rows[0]) if rows else ["item", "value"]
            sheet.append(headers)
            for row in rows:
                sheet.append([row.get(header, "") for header in headers])
            workbook.save(path)
            return DocumentResult(str(path), "xlsx")
        except Exception as exc:
            result = self.create_csv(rows)
            result.fallback = True
            result.metadata = {"requested": "xlsx", "error": str(exc)}
            return result

    def create_csv(self, rows: list[dict[str, Any]]) -> DocumentResult:
        path = self._path("table", "csv")
        headers = list(rows[0]) if rows else ["item", "value"]
        with path.open("w", newline="", encoding="utf-8") as handle:
            writer = csv.DictWriter(handle, fieldnames=headers)
            writer.writeheader()
            writer.writerows(rows)
        return DocumentResult(str(path), "csv")

    def _fallback(self, content: DocumentContent, requested: str, exc: Exception) -> DocumentResult:
        result = self.create_markdown(content)
        result.fallback = True
        result.metadata = {"requested": requested, "error": str(exc)}
        return result

    @staticmethod
    def _rows(content: DocumentContent) -> list[dict[str, Any]]:
        rows = [{"section": "summary", "value": content.summary}]
        for section in content.sections:
            rows.append({"section": section.get("heading", "Section"), "value": section.get("body", "")})
            rows.extend({"section": section.get("heading", "Section"), "value": item} for item in section.get("bullets", []))
        return rows

    def _path(self, title: str, extension: str) -> Path:
        return self.output_dir / f"{_slug(title)}-{time.time_ns()}.{extension}"
